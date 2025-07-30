# Copyright (C) 2024 dssTools Developers
# GNU General Public License v3.0+ (see LICENSE or https://www.gnu.org/licenses/gpl-3.0.txt)
"""Module for handling the batch export of ImageGenerators."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path
from typing import Iterable, Optional

from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.figure import Figure
from pptx import Presentation
from pptx.util import Cm, Length
from pptx import presentation

from dsstools.draw import ImageGenerator, ensure_file_format


class ImageCollection(list):
    """Class for exporting multiple ImageGenerators in one go."""
    def __init__(self, iterable: Optional[Iterable[ImageGenerator]]=None):
        if iterable:
            super().__init__(self._validate_image_generator(item) for item in iterable)

    def __setitem__(self, id, item):
        super().__setitem__(id, self._validate_image_generator(item))

    def append(self, item: ImageGenerator):
        """Add new ImageGenerator to ImageCollection.

        Args:
          item: Item to append to list.

        Returns:

        """
        super().append(self._validate_image_generator(item))

    def extend(self, other: Iterable[ImageGenerator]):
        """Extend existing ImageCollection with another one.

        Args:
          other: Another ImageCollection to extend with.

        Returns:

        """
        if isinstance(other, type(self)):
            super().extend(other)
        else:
            super().extend(self._validate_image_generator(item) for item in other)

    def insert(self, id, item: ImageGenerator):
        """Insert an item at a specific spot.

        Args:
          id: Spot to insert at.
          item: The item to insert.

        Returns:
          The updated ImageCollection.
        """
        super().insert(id, self._validate_image_generator(item))

    def _validate_image_generator(self, value):
        """Ensure only ImageGenerators are being added.

        Args:
          value: Argument to check

        Returns:
          value, if it was actually an ImageGenerator.

        Raises:
          TypeError, if no ImageGenerator was passed.

        """
        if isinstance(value, ImageGenerator):
            return value
        raise TypeError(f"ImageGenerator expected, got {type(value).__name__} instead.")

    def create_flipbook(self, path: Path | str, **kwargs):
        """"Creates a flipbook as PPTX or PDF depending on file ending.

        For the specific valid keyword arguments see `create_flipbook_pdf` or
        `create_flipbook_pptx` which this a wrapper for.

        Args:
            path: Path to save flipbook to. File ending decides on the internal file
                format.

        Returns:
            Either a PDF or PPTX object.
        """
        path = Path(path)
        # TODO Replace this with Henriks file ending checker
        if path.suffix == ".pdf":
            return self.create_flipbook_pdf(path, **kwargs)
        if path.suffix == ".pptx":
            return self.create_flipbook_pptx(path, **kwargs)
        raise ValueError(f"Unable to inpret the file ending {path.suffix}")


    def create_flipbook_pdf(self, path: Path | str) -> PdfPages:
        """Create PDF containing all ImageGenerators.

        Args:
          path: Path to save PDF to.

        Returns:
          Generated PDF object.
        """
        path = Path(path)
        with PdfPages(path) as pdf:
            for ig in self:
                if not ig._fig:
                    ig.draw()
                pdf.savefig(ig._fig)
            return pdf

    def create_flipbook_pptx(
        self,
        path: Path | str,
        titles: Optional[list] = None,
        left: Length = Cm(4),
        top: Length = Cm(-5.3),
        height: Length = Cm(25),
    ) -> presentation.Presentation:
        """Create PPTX containing all ImageGenerators.

        Args:
          path: Path to save file to.
          titles: Titles to give each slide. (Default value = None)
          left:  Left offset of the image on the slide, starting from upper left.
            (Default value = Cm(4))
          top: Top offset of the image on the slide, starting from upper left.
            (Default value = Cm(-5.3))
          height: Height of the image. By default uses a sensible default. If you change
            this, you might have to adapt the left and top arguments as well. (Default
            value = Cm(25))

        Returns:
          Generated PPTX object.
        """
        path = Path(path)

        if titles and len(self) != len(titles):
            raise ValueError(
                "Titles and image generators in the flipbook are not the same size."
            )

        prs = Presentation()
        # We default to 16:9 as slide layout
        prs.slide_width = 9144000
        prs.slide_height = 5144400

        for i, ig in enumerate(self):
            title_slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(title_slide_layout)
            if titles:
                slide.shapes.title.text = titles[i]
            if not ig._fig:
                ig.draw()
            # Save the image to a byte stream to avoid disk access
            stream = BytesIO()
            ig._fig.savefig(stream, dpi=ig.dpi)
            pic = slide.shapes.add_picture(stream, left=left, top=top, height=height)

            # Move image to background
            cursor_sp = slide.shapes[0]._element
            cursor_sp.addprevious(pic._element)

        prs.save(path)
        return prs

    def create_multiple_in_one(self, fig: Figure, path: Path | str, dpi: int = 200):
        path = Path(path)

        # TODO Use Henriks file ending checker
        if len(self) != len(fig.axes):
            raise ValueError("List of graphs and axes in figure are not the same length.")

        for ig, ax in zip(self, fig.axes):
            ig.set_axis(ax).draw()
            ax.set_axis_off()

        if path:
            fpath, file_format = ensure_file_format(path,
                                                    path.suffix,
                                                    default_format=".svg")
            fig.savefig(fpath, format=file_format, dpi=dpi)

        return fig
