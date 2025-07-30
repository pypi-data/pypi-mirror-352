from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd
import os, asyncio
from .multitask import run_in_thread_pool

# for now, only return 1 page
def load_word(filepath_:dir)->list[str]:
    paragraghs = []
    document = Document(filepath_)
    for p in document.paragraphs:
        if p.text and p.text.strip():
            paragraghs.append(p.text)
    return ['\n'.join(paragraghs)]

def load_pdf(filepath_:dir)->list[str]:
    reader = PdfReader(filepath_)
    return [page.extract_text() for page in reader.pages]

def load_excel(filepath_:dir)->list[str]:
    xls = pd.ExcelFile(filepath_)
    pages = []
    for sheet in xls.sheet_names:
        df = pd.read_excel(filepath_, sheet_name=sheet)
        pages.append(str(df))
    return  pages  

def load_text(filepath_:dir)->list[str]:
    with open(filepath_, "r", encoding='utf8') as f:
        content = f.read()
        return [content] if content else []

def load_pptx(filepath_:dir)->list[str]:
    prs = Presentation(filepath_)
    pages = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    texts.append(run.text)
        pages.append('\n'.join(texts))
    return pages

## TO-DO: add other file types
async def load_document_pages(filepath_:dir)->list:
    return await run_in_thread_pool(_load_document_local, filepath_)

async def load_document_content(filepath_:dir)->str:
    pages = await load_document_pages(filepath_)
    return '\n\n'.join(pages)

def _load_document_local(filepath_:dir)->list[str]:
    _, ext = os.path.splitext(filepath_)
    ext = ext.lower()
    if ext.endswith("pdf"):
        return load_pdf(filepath_)
    elif ext.endswith(("doc", "docx")):
        return load_word(filepath_)
    elif ext.endswith(("xls", "xlsx", "xlsd")):
        return load_excel(filepath_)
    elif ext.endswith(("ppt", "pptx")):
        return load_pptx(filepath_)
    else:
        return load_text(filepath_)
    
def _load_document_web(url_:str)->list[str]:
    pass
# convert word to pdf to read pages
"""
import docx2pdf
from PyPDF2 import PdfReader
p = docx2pdf.convert('myword.docx','document.pdf')
r = PdfReader('document.pdf')
num_pages = len(r.pages)
"""

async def main():
    while True:
        input_file = input("Input file to save (or Enter to skip):")
        input_file = input_file.strip()
        if not input_file:
            break
        input_file = input_file.replace('"', '')
        content = await load_document_content(input_file)
        print(content)

if __name__ == '__main__':
    asyncio.run(main())