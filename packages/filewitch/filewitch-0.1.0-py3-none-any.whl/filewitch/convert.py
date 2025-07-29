"""
Core conversion functionality for FileWitch.
"""

import pandas as pd
import openpyxl
from pathlib import Path
import logging
import docx
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
import shutil
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from pptx import Presentation

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class ConversionError(Exception):
    """Custom exception for conversion errors."""
    pass

def csv_to_xlsx(input_path: str, output_path: str) -> None:
    """Convert CSV file to Excel format."""
    try:
        df = pd.read_csv(input_path)
        wb = openpyxl.Workbook()
        ws = wb.active
        
        # Write headers
        for col_num, column in enumerate(df.columns, 1):
            ws.cell(row=1, column=col_num, value=column)
        
        # Write data
        for row_num, row in enumerate(df.values, 2):
            for col_num, value in enumerate(row, 1):
                ws.cell(row=row_num, column=col_num, value=value)
        
        wb.save(output_path)
        logger.info(f"Successfully converted {input_path} to {output_path}")
    except Exception as e:
        raise ConversionError(f"Failed to convert CSV to Excel: {str(e)}")

def xlsx_to_csv(input_path: str, output_path: str) -> None:
    """Convert Excel file to CSV format."""
    try:
        df = pd.read_excel(input_path)
        df.to_csv(output_path, index=False)
        logger.info(f"Successfully converted {input_path} to {output_path}")
    except Exception as e:
        raise ConversionError(f"Failed to convert Excel to CSV: {str(e)}")

def txt_to_docx(input_path: str, output_path: str) -> None:
    """Convert text file to Word document."""
    try:
        with open(input_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        doc = docx.Document()
        paragraph = doc.add_paragraph()
        paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
        run = paragraph.add_run(content)
        run.font.size = Pt(11)
        run.font.name = 'Arial'
        
        doc.save(output_path)
        logger.info(f"Successfully converted {input_path} to {output_path}")
    except Exception as e:
        raise ConversionError(f"Failed to convert text to Word: {str(e)}")

def docx_to_txt(input_path: str, output_path: str) -> None:
    """Convert Word document to text file."""
    try:
        doc = docx.Document(input_path)
        text_content = [paragraph.text for paragraph in doc.paragraphs]
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(text_content))
        
        logger.info(f"Successfully converted {input_path} to {output_path}")
    except Exception as e:
        raise ConversionError(f"Failed to convert Word to text: {str(e)}")

def copy_file(input_path: str, output_path: str) -> None:
    """Copy a file from one location to another."""
    try:
        shutil.copy2(input_path, output_path)
        logger.info(f"Successfully copied {input_path} to {output_path}")
    except Exception as e:
        raise ConversionError(f"Failed to copy file: {str(e)}")

def txt_to_pdf(input_path: str, output_path: str) -> None:
    """Convert text file to PDF format."""
    try:
        with open(input_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        
        # Set font and size
        c.setFont("Helvetica", 12)
        
        # Split content into lines and write to PDF
        y = height - 50  # Start from top with margin
        for line in content.split('\n'):
            if y < 50:  # If we're near bottom, start new page
                c.showPage()
                c.setFont("Helvetica", 12)
                y = height - 50
            c.drawString(50, y, line)
            y -= 15  # Move down for next line
        
        c.save()
        logger.info(f"Successfully converted {input_path} to {output_path}")
    except Exception as e:
        raise ConversionError(f"Failed to convert text to PDF: {str(e)}")

def docx_to_pdf(input_path: str, output_path: str) -> None:
    """Convert Word document to PDF format."""
    try:
        doc = docx.Document(input_path)
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        
        # Set font and size
        c.setFont("Helvetica", 12)
        
        # Write content to PDF
        y = height - 50  # Start from top with margin
        for paragraph in doc.paragraphs:
            if y < 50:  # If we're near bottom, start new page
                c.showPage()
                c.setFont("Helvetica", 12)
                y = height - 50
            c.drawString(50, y, paragraph.text)
            y -= 15  # Move down for next line
        
        c.save()
        logger.info(f"Successfully converted {input_path} to {output_path}")
    except Exception as e:
        raise ConversionError(f"Failed to convert Word to PDF: {str(e)}")

def pptx_to_docx(input_path: str, output_path: str) -> None:
    """Convert PowerPoint presentation to Word document."""
    try:
        prs = Presentation(input_path)
        doc = docx.Document()
        
        # Process each slide
        for slide in prs.slides:
            # Add slide title
            if slide.shapes.title:
                doc.add_heading(slide.shapes.title.text, level=1)
            
            # Add slide content
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                    doc.add_paragraph(shape.text)
            
            # Add a page break between slides
            doc.add_page_break()
        
        doc.save(output_path)
        logger.info(f"Successfully converted {input_path} to {output_path}")
    except Exception as e:
        raise ConversionError(f"Failed to convert PowerPoint to Word: {str(e)}")