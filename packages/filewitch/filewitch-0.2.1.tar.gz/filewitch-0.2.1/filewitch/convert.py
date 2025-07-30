"""
Core conversion functionality for FileWitch.
"""

import pandas as pd
import openpyxl
from pathlib import Path
import logging
import docx
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
import shutil
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from pptx import Presentation
from PIL import Image
import io
import os
from docx2pdf import convert as docx2pdf_convert
# from pptx2pdf import convert as pptx2pdf_convert  # Removed because pptx2pdf does not exist

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class ConversionError(Exception):
    """Custom exception for conversion errors."""
    pass

def _extract_images_from_docx(doc):
    """Extract images from a Word document."""
    images = []
    for rel in doc.part.rels.values():
        if "image" in rel.target_ref:
            image_data = rel.target_part.blob
            images.append(image_data)
    return images

def _add_image_to_docx(doc, image_data, width=Inches(6)):
    """Add an image to a Word document."""
    image_stream = io.BytesIO(image_data)
    doc.add_picture(image_stream, width=width)

def csv_to_xlsx(input_path: str, output_path: str) -> None:
    """Convert CSV file to Excel format."""
    try:
        df = pd.read_csv(input_path)
        wb = openpyxl.Workbook()
        ws = wb.active
        
        # Write headers
        for col_num, column in enumerate(df.columns, 1):
            ws.cell(row=1, column=col_num, value=column)
        
        # Write data
        for row_num, row in enumerate(df.values, 2):
            for col_num, value in enumerate(row, 1):
                ws.cell(row=row_num, column=col_num, value=value)
        
        wb.save(output_path)
        logger.info(f"Successfully converted {input_path} to {output_path}")
    except Exception as e:
        raise ConversionError(f"Failed to convert CSV to Excel: {str(e)}")

def xlsx_to_csv(input_path: str, output_path: str) -> None:
    """Convert Excel file to CSV format."""
    try:
        df = pd.read_excel(input_path)
        df.to_csv(output_path, index=False)
        logger.info(f"Successfully converted {input_path} to {output_path}")
    except Exception as e:
        raise ConversionError(f"Failed to convert Excel to CSV: {str(e)}")

def txt_to_docx(input_path: str, output_path: str) -> None:
    """Convert text file to Word document with basic formatting."""
    try:
        with open(input_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        doc = docx.Document()
        
        # Process content line by line
        for line in content.split('\n'):
            if line.strip() == '[IMAGE]':
                # Skip image placeholders
                continue
            elif line.strip().startswith('#'):
                # Handle headings
                level = len(line.split()[0])
                doc.add_heading(line.lstrip('#').strip(), level=min(level, 9))
            else:
                # Regular paragraph
                paragraph = doc.add_paragraph()
                paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT
                run = paragraph.add_run(line)
                run.font.size = Pt(11)
                run.font.name = 'Arial'
        
        doc.save(output_path)
        logger.info(f"Successfully converted {input_path} to {output_path}")
    except Exception as e:
        raise ConversionError(f"Failed to convert text to Word: {str(e)}")

def docx_to_txt(input_path: str, output_path: str) -> None:
    """Convert Word document to text file with image placeholders."""
    try:
        doc = docx.Document(input_path)
        text_content = []
        
        for paragraph in doc.paragraphs:
            # Add paragraph text
            text_content.append(paragraph.text)
            
            # Check for images in the paragraph
            for run in paragraph.runs:
                if run._element.findall('.//w:drawing'):
                    text_content.append("[IMAGE]")
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(text_content))
        
        logger.info(f"Successfully converted {input_path} to {output_path}")
    except Exception as e:
        raise ConversionError(f"Failed to convert Word to text: {str(e)}")

def copy_file(input_path: str, output_path: str) -> None:
    """Copy a file from one location to another."""
    try:
        shutil.copy2(input_path, output_path)
        logger.info(f"Successfully copied {input_path} to {output_path}")
    except Exception as e:
        raise ConversionError(f"Failed to copy file: {str(e)}")

def txt_to_pdf(input_path: str, output_path: str) -> None:
    """Convert text file to PDF format."""
    try:
        with open(input_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        c = canvas.Canvas(output_path, pagesize=letter)
        width, height = letter
        
        # Set font and size
        c.setFont("Helvetica", 12)
        
        # Split content into lines and write to PDF
        y = height - 50  # Start from top with margin
        for line in content.split('\n'):
            if y < 50:  # If we're near bottom, start new page
                c.showPage()
                c.setFont("Helvetica", 12)
                y = height - 50
            c.drawString(50, y, line)
            y -= 15  # Move down for next line
        
        c.save()
        logger.info(f"Successfully converted {input_path} to {output_path}")
    except Exception as e:
        raise ConversionError(f"Failed to convert text to PDF: {str(e)}")

def docx_to_pdf(input_path: str, output_path: str) -> None:
    """Convert Word document to PDF format with preserved formatting."""
    try:
        # Use docx2pdf for better formatting preservation
        docx2pdf_convert(input_path, output_path)
        logger.info(f"Successfully converted {input_path} to {output_path}")
    except Exception as e:
        raise ConversionError(f"Failed to convert Word to PDF: {str(e)}")

def pptx_to_docx(input_path: str, output_path: str) -> None:
    """Convert PowerPoint presentation to Word document with images and formatting."""
    try:
        prs = Presentation(input_path)
        doc = docx.Document()
        
        # Process each slide
        for slide in prs.slides:
            # Add slide title
            if slide.shapes.title:
                doc.add_heading(slide.shapes.title.text, level=1)
            
            # Process shapes (text and images)
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape != slide.shapes.title:
                    # Add text content
                    doc.add_paragraph(shape.text)
                
                # Handle images
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        image_stream = io.BytesIO(shape.image.blob)
                        doc.add_picture(image_stream, width=Inches(6))
                    except Exception as e:
                        logger.warning(f"Failed to process image in slide: {str(e)}")
            
            # Add a page break between slides
            doc.add_page_break()
        
        doc.save(output_path)
        logger.info(f"Successfully converted {input_path} to {output_path}")
    except Exception as e:
        raise ConversionError(f"Failed to convert PowerPoint to Word: {str(e)}")

def pptx_to_pdf(input_path: str, output_path: str) -> None:
    """Convert PowerPoint presentation to PDF with preserved formatting. (Feature disabled: pptx2pdf not available)"""
    raise NotImplementedError("pptx2pdf is not available. This feature is currently disabled.")