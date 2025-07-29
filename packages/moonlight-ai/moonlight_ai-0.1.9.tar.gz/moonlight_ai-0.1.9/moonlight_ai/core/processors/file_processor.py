import sqlite3, docx, PyPDF2, csv
from pptx import Presentation
from pathlib import Path

from .base import ProcessorException

class FileProcessor:
    def __init__(self, files: list = []):
        self.files = files
        self.content = {}
        for file in files:
            if not Path(file).exists():
                raise ProcessorException(f"File not found: {file}")
            self.content[file] = self.process(file)
        
    def process(self, file_path: str) -> str:
        file = Path(file_path)
        ext = file.suffix.lower()
        if ext == ".txt":
            return self._process_txt(file_path)
        elif ext == ".docx":
            return self._process_docx(file_path)
        elif ext == ".pdf":
            return self._process_pdf(file_path)
        elif ext == ".csv":
            return self._process_csv(file_path)
        elif ext in [".xlsx", ".xls"]:
            return self._process_excel(file_path)
        elif ext in [".db", ".sqlite", ".sqlite3"]:
            return self._process_db(file_path)
        elif ext in [".pptx", ".ppt"]:
            return self._process_pptx(file_path)
        else:
            raise ProcessorException(f"Unsupported file type: {ext}")

    def _process_txt(self, file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    def _process_docx(self, file_path: str) -> str:
        doc = docx.Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text]
        return "\n".join(paragraphs)

    def _process_pdf(self, file_path: str) -> str:
        text = ""
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
        return text

    def _process_csv(self, file_path: str) -> str:
        with open(file_path, newline="", encoding="utf-8") as csvfile:
            reader = csv.reader(csvfile)
            rows = ["\t".join(row) for row in reader]
        return "\n".join(rows)

    def _process_excel(self, file_path: str) -> str:
        ext = Path(file_path).suffix.lower()
        texts = []
        if ext == ".xls":
            try:
                import xlrd
            except ImportError:
                raise ProcessorException("xlrd module is required for .xls files.")
            try:
                workbook = xlrd.open_workbook(file_path)
            except xlrd.compdoc.CompDocError as e:
                raise ProcessorException(
                    f"Failed to process .xls file {file_path}: {str(e)}. Please convert it to .xlsx or check for corruption."
                )
            for sheet in workbook.sheets():
                for rownum in range(sheet.nrows):
                    row = sheet.row_values(rownum)
                    texts.append("\t".join([str(cell) for cell in row]))
        else:
            try:
                import openpyxl
            except ImportError:
                raise ProcessorException("openpyxl module is required for Excel files.")
            wb = openpyxl.load_workbook(file_path, data_only=True)
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    texts.append("\t".join([str(cell) if cell is not None else "" for cell in row]))
        return "\n".join(texts)

    def _process_db(self, file_path: str) -> str:
        conn = sqlite3.connect(file_path)
        cursor = conn.cursor()
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
        table_names = [table[0] for table in cursor.fetchall()]
        if not table_names:
            conn.close()
            return "No tables found."
        content = []
        for table in table_names:
            content.append(f"Table: {table}")
            # Get column names
            cursor.execute(f"PRAGMA table_info('{table}')")
            columns = [col[1] for col in cursor.fetchall()]
            content.append("\t".join(columns))
            # Get all rows from table
            cursor.execute(f"SELECT * FROM \"{table}\"")
            rows = cursor.fetchall()
            for row in rows:
                content.append("\t".join(str(item) for item in row))
            content.append("")  # Add an empty line between tables
        conn.close()
        return "\n".join(content)
    
    def _process_pptx(self, file_path: str) -> str:        
        texts = []
        prs = Presentation(file_path)
        
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            texts.append("\n".join(slide_text))
        
        return "\n\n".join(texts)