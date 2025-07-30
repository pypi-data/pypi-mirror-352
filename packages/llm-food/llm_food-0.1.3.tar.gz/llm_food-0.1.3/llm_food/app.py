import asyncio
import base64
import shutil
import tempfile

from datetime import datetime
import hashlib
import json
from io import BytesIO
import os
from typing import List, Optional, Dict
import uuid
from fastapi import (
    FastAPI,
    UploadFile,
    File,
    BackgroundTasks,
    HTTPException,
    Query,
    Depends,
    Form,
)
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials

import httpx
import mammoth
import duckdb

# Imports for GCS
from google.cloud import storage
from google.oauth2 import service_account  # For local testing with service account
from markdownify import markdownify

from striprtf.striprtf import rtf_to_text
from pptx import Presentation
import trafilatura

from .config import (
    get_pdf_backend,
    get_gemini_prompt,
    get_api_auth_token,
    get_gcs_project_id,
    get_max_file_size_bytes,
    DUCKDB_FILE,
    GCS_BATCH_BUCKET,
    GEMINI_MODEL_FOR_VISION,
    SUPPORTED_EXTENSIONS,
)

from .models import (
    BatchJobFailedFileOutput,
    BatchJobOutputResponse,
    BatchOutputItem,
    ConversionResponse,
    BatchJobStatusResponse,
)

# --- Conditional imports based on the PDF backend ---
match get_pdf_backend():
    case "pymupdf4llm":
        from pymupdf4llm import to_markdown
        import pymupdf
    case "pypdf2":
        from pypdf import PdfReader
    case "gemini":
        from pdf2image import convert_from_bytes
        from google import genai
        from google.genai.types import CreateBatchJobConfig, JobState

        OCR_PROMPT = get_gemini_prompt()
    case invalid_backend:
        raise ValueError(f"Invalid PDF backend: {invalid_backend}")

# --- Main application setup ---
app = FastAPI(
    title="LLM Food API",
    description="API for converting various document formats to Markdown or text, with batch processing capabilities.",
    version="0.2.0",
)

# --- Security ---
bearer_scheme = HTTPBearer(
    auto_error=False
)  # auto_error=False to handle optional token & custom errors


async def authenticate_request(
    credentials: Optional[HTTPAuthorizationCredentials] = Depends(bearer_scheme),
) -> None:
    configured_token = get_api_auth_token()
    if configured_token:  # Only enforce auth if a token is configured server-side
        if credentials is None:
            raise HTTPException(
                status_code=401,
                detail="Not authenticated. Authorization header is missing.",
                headers={"WWW-Authenticate": "Bearer"},
            )
        if credentials.scheme != "Bearer":
            raise HTTPException(
                status_code=401,
                detail="Invalid authentication scheme. Only Bearer is supported.",
                headers={"WWW-Authenticate": "Bearer"},
            )
        if credentials.credentials != configured_token:
            raise HTTPException(
                status_code=403,
                detail="Invalid token.",
            )
    # If no token is configured server-side, or if authentication passes, do nothing.
    return


# For local GCS testing with a service account JSON file
def get_gcs_credentials():
    credentials_path = os.getenv("GOOGLE_APPLICATION_CREDENTIALS")
    if credentials_path:
        return service_account.Credentials.from_service_account_file(credentials_path)
    return None  # Fallback to default environment auth if not set


def get_gemini_client():
    project = get_gcs_project_id()
    location = os.getenv("GOOGLE_CLOUD_LOCATION", "us-central1")
    client = genai.Client(vertexai=True, location=location, project=project)
    return client


TASKS = {}


def get_db_connection():
    return duckdb.connect(DUCKDB_FILE)


def initialize_db_schema():
    con = get_db_connection()
    try:
        # Main batch jobs table
        con.execute("""
            CREATE TABLE IF NOT EXISTS batch_jobs (
                job_id VARCHAR PRIMARY KEY,
                output_gcs_path VARCHAR NOT NULL,
                status VARCHAR NOT NULL,
                submitted_at TIMESTAMP NOT NULL,
                total_input_files INTEGER NOT NULL,
                overall_processed_count INTEGER DEFAULT 0,
                overall_failed_count INTEGER DEFAULT 0,
                last_updated_at TIMESTAMP
            )
        """)
        # Gemini PDF sub-jobs (one per Gemini Batch API call)
        con.execute("""
            CREATE TABLE IF NOT EXISTS gemini_pdf_batch_sub_jobs (
                gemini_batch_sub_job_id VARCHAR PRIMARY KEY,
                batch_job_id VARCHAR NOT NULL REFERENCES batch_jobs(job_id),
                gemini_api_job_name VARCHAR,
                status VARCHAR NOT NULL,
                payload_gcs_uri VARCHAR,
                gemini_output_gcs_uri_prefix VARCHAR,
                total_pdf_pages_for_batch INTEGER DEFAULT 0,
                processed_pdf_pages_count INTEGER DEFAULT 0,
                failed_pdf_pages_count INTEGER DEFAULT 0,
                error_message VARCHAR,
                created_at TIMESTAMP NOT NULL,
                updated_at TIMESTAMP
            )
        """)
        # Individual file tasks (for non-PDFs, or individual pages of PDFs before aggregation)
        con.execute("""
            CREATE TABLE IF NOT EXISTS file_tasks (
                file_task_id VARCHAR PRIMARY KEY,
                batch_job_id VARCHAR NOT NULL REFERENCES batch_jobs(job_id),
                gemini_batch_sub_job_id VARCHAR REFERENCES gemini_pdf_batch_sub_jobs(gemini_batch_sub_job_id), -- Link to a Gemini batch if it's a PDF page
                original_filename VARCHAR NOT NULL,
                file_type VARCHAR NOT NULL, -- e.g., 'pdf_page', 'docx'
                status VARCHAR NOT NULL, -- pending, processing, image_uploaded_to_gcs, completed, failed
                gcs_input_image_uri VARCHAR, -- For PDF pages, GCS URI of the image sent to Gemini
                gcs_output_markdown_uri VARCHAR, -- GCS URI of the final .md (for non-PDFs or aggregated PDFs)
                page_number INTEGER, -- For PDF pages
                gemini_request_id VARCHAR, -- The 'id' used in payload.jsonl for this PDF page
                error_message VARCHAR,
                created_at TIMESTAMP NOT NULL,
                updated_at TIMESTAMP
            )
        """)
    finally:
        con.close()


# Call initialization at startup
initialize_db_schema()


def _process_docx_sync(content_bytes: bytes) -> List[str]:
    try:
        doc = BytesIO(content_bytes)
        doc_html = mammoth.convert_to_html(doc).value
        doc_md = markdownify(doc_html).strip()
        return [doc_md]
    except Exception as e:
        return [f"Error processing DOCX: {str(e)}"]


def _process_rtf_sync(content_bytes: bytes) -> List[str]:
    try:
        return [rtf_to_text(content_bytes.decode("utf-8", errors="ignore"))]
    except Exception as e:
        return [f"Error processing RTF: {str(e)}"]


def _process_pptx_sync(content_bytes: bytes) -> List[str]:
    try:
        prs = Presentation(BytesIO(content_bytes))
        # Corrected list comprehension for PPTX to build a single string per slide, then list of slide texts
        slide_texts = []
        for slide in prs.slides:
            text_on_slide = "\n".join(
                shape.text
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text
            )
            if text_on_slide:  # Only add if there's text
                slide_texts.append(text_on_slide)
        return (
            slide_texts if slide_texts else [""]
        )  # Return list of slide texts, or list with empty string if no text
    except Exception as e:
        return [f"Error processing PPTX: {str(e)}"]


def _process_html_sync(content_bytes: bytes) -> List[str]:
    try:
        extracted_text = trafilatura.extract(
            content_bytes.decode("utf-8", errors="ignore"), output_format="markdown"
        )
        return [extracted_text if extracted_text is not None else ""]
    except Exception as e:
        return [f"Error processing HTML: {str(e)}"]


def _process_pdf_pymupdf4llm_sync(content_bytes: bytes) -> List[str]:
    try:
        pymupdf_doc = pymupdf.Document(stream=content_bytes, filetype="pdf")
        page_data_list = to_markdown(pymupdf_doc, page_chunks=True)
        return [page_dict.get("text", "") for page_dict in page_data_list]
    except Exception as e:
        return [f"Error processing PDF with pymupdf4llm: {str(e)}"]


def _process_pdf_pypdf2_sync(content_bytes: bytes) -> List[str]:
    try:
        reader = PdfReader(BytesIO(content_bytes))
        return [p.extract_text() or "" for p in reader.pages]
    except Exception as e:
        return [f"Error processing PDF with pypdf: {str(e)}"]


async def _process_file_content(
    ext: str, content: bytes, pdf_backend_choice: str
) -> List[str]:
    texts_list: List[str] = []
    if ext == ".pdf":
        if pdf_backend_choice == "pymupdf4llm":
            texts_list = await asyncio.to_thread(_process_pdf_pymupdf4llm_sync, content)
        elif pdf_backend_choice == "pypdf2":
            texts_list = await asyncio.to_thread(_process_pdf_pypdf2_sync, content)
        elif pdf_backend_choice == "gemini":
            pages = convert_from_bytes(content)
            images_b64 = []
            for page in pages:
                buffer = BytesIO()
                page.save(buffer, format="PNG")
                image_data = buffer.getvalue()
                b64_str = base64.b64encode(image_data).decode("utf-8")
                images_b64.append(b64_str)
            client = get_gemini_client()
            payloads = [
                [
                    {"inline_data": {"data": b64_str, "mime_type": "image/png"}},
                    {"text": OCR_PROMPT},
                ]
                for b64_str in images_b64
            ]
            results = await asyncio.gather(
                *[
                    client.aio.models.generate_content(
                        model=GEMINI_MODEL_FOR_VISION, contents=payload
                    )
                    for payload in payloads
                ]
            )
            texts_list = [result.text for result in results]
        else:
            texts_list = ["Invalid PDF backend specified."]
    elif ext in [".docx"]:
        texts_list = await asyncio.to_thread(_process_docx_sync, content)
    elif ext in [".rtf"]:
        texts_list = await asyncio.to_thread(_process_rtf_sync, content)
    elif ext in [".pptx"]:
        texts_list = await asyncio.to_thread(_process_pptx_sync, content)
    elif ext in [".html", ".htm"]:
        texts_list = await asyncio.to_thread(_process_html_sync, content)
    else:
        texts_list = ["Unsupported file type encountered in _process_file_content."]
    return texts_list


@app.post(
    "/convert",
    response_model=ConversionResponse,
    dependencies=[Depends(authenticate_request)],
)
async def convert_file_upload(file: UploadFile = File(...)):
    ext = os.path.splitext(file.filename)[1].lower()
    content = await file.read()

    max_size = get_max_file_size_bytes()
    if max_size is not None and len(content) > max_size:
        raise HTTPException(
            status_code=413,
            detail=f"File size {len(content) / (1024 * 1024):.2f}MB exceeds maximum allowed size of {max_size / (1024 * 1024):.2f}MB.",
        )

    content_hash = hashlib.sha256(content).hexdigest()
    pdf_backend_choice = get_pdf_backend()

    texts_list = await _process_file_content(ext, content, pdf_backend_choice)

    if texts_list and (
        texts_list[0].startswith("Error processing")
        or texts_list[0].startswith("Invalid PDF backend")
        or texts_list[0].startswith("Unsupported file type")
    ):
        raise HTTPException(status_code=400, detail=texts_list[0])

    return ConversionResponse(
        filename=file.filename, content_hash=content_hash, texts=texts_list
    )


@app.get(
    "/convert",
    response_model=ConversionResponse,
    dependencies=[Depends(authenticate_request)],
)
async def convert_url(
    url: str = Query(..., description="URL of the webpage to convert to Markdown"),
):
    try:
        async with httpx.AsyncClient() as client:
            response = await client.get(url)
            response.raise_for_status()
            html_content = response.text
            content_bytes = html_content.encode("utf-8")
    except httpx.RequestError as e:
        raise HTTPException(status_code=400, detail=f"Error fetching URL: {str(e)}")
    except httpx.HTTPStatusError as e:
        raise HTTPException(
            status_code=e.response.status_code,
            detail=f"Error fetching URL: {e.response.reason_phrase}",
        )

    if not html_content:
        raise HTTPException(status_code=400, detail="Fetched content is empty.")

    content_hash = hashlib.sha256(content_bytes).hexdigest()

    extracted_text = trafilatura.extract(html_content, output_format="markdown")
    texts_list = [extracted_text if extracted_text is not None else ""]

    filename = os.path.basename(url) or url

    return ConversionResponse(
        filename=filename, content_hash=content_hash, texts=texts_list
    )


@app.get(
    "/status/{task_id}",
    response_model=BatchJobStatusResponse,
    dependencies=[Depends(authenticate_request)],
)
def status(task_id: str):
    con = get_db_connection()
    try:
        job_status_row = con.execute(
            "SELECT * FROM batch_jobs WHERE job_id = ?", (task_id,)
        ).fetchone()

        if not job_status_row:
            raise HTTPException(
                status_code=404, detail=f"Batch job with ID {task_id} not found."
            )

        job_dict = dict(zip([desc[0] for desc in con.description], job_status_row))

        gemini_sub_jobs_rows = con.execute(
            "SELECT * FROM gemini_pdf_batch_sub_jobs WHERE batch_job_id = ?",
            (task_id,),
        ).fetchall()
        job_dict["gemini_pdf_processing_details"] = [
            dict(zip([desc[0] for desc in con.description], sub_job_row))
            for sub_job_row in gemini_sub_jobs_rows
        ]

        file_tasks_rows = con.execute(
            "SELECT original_filename, file_type, status, gcs_output_markdown_uri, error_message, page_number FROM file_tasks WHERE batch_job_id = ?",
            (task_id,),
        ).fetchall()
        job_dict["file_processing_details"] = [
            dict(zip([desc[0] for desc in con.description], task_row))
            for task_row in file_tasks_rows
        ]

        # Pydantic will validate the structure of job_dict against BatchJobStatusResponse
        return job_dict
    finally:
        con.close()


@app.post("/batch", dependencies=[Depends(authenticate_request)])
def batch_files_upload(
    background_tasks: BackgroundTasks,
    files: List[UploadFile] = File(...),
    output_gcs_path: str = Form(...),
):
    main_batch_job_id = str(uuid.uuid4())
    current_time = datetime.utcnow()

    pdf_files_data_for_batch: List[
        tuple[str, str]
    ] = []  # Changed: Store (filename, temp_file_path)
    non_pdf_files_for_individual_processing: List[
        tuple[str, str, str]
    ] = []  # Store (filename, ext, temp_file_path)

    if not GCS_BATCH_BUCKET:
        raise HTTPException(
            status_code=500,
            detail="GCS_BATCH_TEMP_BUCKET is not configured on the server.",
        )
    if not output_gcs_path.startswith("gs://"):
        raise HTTPException(
            status_code=400, detail="Output GCS path must start with gs://"
        )

    temp_files = []  # Keep track of temp files for cleanup in case of errors
    try:
        print("Processing uploaded files...")
        for f in files:
            ext = os.path.splitext(f.filename)[1].lower()

            # Create a temporary file with the original extension
            temp_file = tempfile.NamedTemporaryFile(suffix=ext, delete=False)
            temp_files.append(temp_file.name)

            # Write uploaded content to temp file
            shutil.copyfileobj(f.file, temp_file)
            temp_file.close()  # Close but don't delete (delete=False above)

            if ext == ".pdf":
                pdf_files_data_for_batch.append((f.filename, temp_file.name))
            elif ext in SUPPORTED_EXTENSIONS:  # Excludes .pdf as it's handled above
                non_pdf_files_for_individual_processing.append(
                    (f.filename, ext, temp_file.name)
                )
            else:
                # Clean up this temp file immediately if unsupported
                os.unlink(temp_file.name)
                temp_files.remove(temp_file.name)
                print(f"Skipping unsupported file: {f.filename}")

        con = get_db_connection()
        try:
            con.execute(
                "INSERT INTO batch_jobs (job_id, output_gcs_path, status, submitted_at, total_input_files, last_updated_at) VALUES (?, ?, ?, ?, ?, ?)",
                (
                    main_batch_job_id,
                    output_gcs_path,
                    "pending",
                    current_time,
                    len(pdf_files_data_for_batch)
                    + len(non_pdf_files_for_individual_processing),
                    current_time,
                ),
            )
            con.commit()

            # Process non-PDF files
            for (
                orig_filename,
                file_ext,
                temp_file_path,
            ) in non_pdf_files_for_individual_processing:
                file_task_id = str(uuid.uuid4())
                con.execute(
                    "INSERT INTO file_tasks (file_task_id, batch_job_id, original_filename, file_type, status, created_at, updated_at) VALUES (?, ?, ?, ?, ?, ?, ?)",
                    (
                        file_task_id,
                        main_batch_job_id,
                        orig_filename,
                        file_ext,
                        "pending",
                        current_time,
                        current_time,
                    ),
                )

                background_tasks.add_task(
                    _process_single_non_pdf_file_and_upload,
                    temp_file_path,
                    file_ext,
                    orig_filename,
                    output_gcs_path,
                    main_batch_job_id,
                    file_task_id,
                )
            con.commit()

            # Process PDF files via Gemini Batch
            if pdf_files_data_for_batch:
                gemini_batch_sub_job_id = str(uuid.uuid4())
                con.execute(
                    "INSERT INTO gemini_pdf_batch_sub_jobs (gemini_batch_sub_job_id, batch_job_id, status, created_at, updated_at) VALUES (?, ?, ?, ?, ?)",
                    (
                        gemini_batch_sub_job_id,
                        main_batch_job_id,
                        "pending_preparation",
                        current_time,
                        current_time,
                    ),
                )
                # Pass the list of (filename, temp_file_path)
                background_tasks.add_task(
                    _run_gemini_pdf_batch_conversion,
                    pdf_files_data_for_batch,
                    output_gcs_path,
                    main_batch_job_id,
                    gemini_batch_sub_job_id,
                )
                print("Added background tasks for batch prediction PDF")
                con.commit()

            # Update batch job status to processing if there are tasks
            if pdf_files_data_for_batch or non_pdf_files_for_individual_processing:
                con.execute(
                    "UPDATE batch_jobs SET status = ?, last_updated_at = ? WHERE job_id = ?",
                    ("processing", datetime.utcnow(), main_batch_job_id),
                )
            else:  # No files to process
                con.execute(
                    "UPDATE batch_jobs SET status = ?, last_updated_at = ? WHERE job_id = ?",
                    ("completed_no_files", datetime.utcnow(), main_batch_job_id),
                )
            con.commit()

        finally:
            _check_and_finalize_batch_job_status(main_batch_job_id, con)
            con.commit()
            con.close()

        return {"task_id": main_batch_job_id}

    except Exception as e:
        # Clean up temp files in case of error
        for temp_file_path in temp_files:
            try:
                os.unlink(temp_file_path)
            except Exception as cleanup_error:
                print(f"Error cleaning up temp file {temp_file_path}: {cleanup_error}")
        raise HTTPException(
            status_code=500,
            detail=f"Error processing batch upload: {str(e)}",
        )


@app.get(
    "/batch/{task_id}",
    response_model=BatchJobOutputResponse,
    dependencies=[Depends(authenticate_request)],
)
async def get_batch_output(task_id: str):
    con = get_db_connection()
    try:
        job_details_tuple = con.execute(
            "SELECT job_id, status, output_gcs_path FROM batch_jobs WHERE job_id = ?",
            (task_id,),
        ).fetchone()
        if not job_details_tuple:
            raise HTTPException(
                status_code=404, detail=f"Batch job {task_id} not found."
            )

        job_id, job_status, output_gcs_path = job_details_tuple

        outputs_list = []
        errors_list = []
        message = None

        completed_statuses_for_output = ["completed", "completed_with_errors"]

        if job_status in completed_statuses_for_output:
            # Fetch distinct successful outputs
            # A file is considered successfully processed if at least one of its file_tasks (e.g. a page for a PDF, or the file itself for docx)
            # has a gcs_output_markdown_uri and status completed.
            # We need the original_filename and the gcs_output_markdown_uri of the *final aggregated file*.
            # The current DB schema stores the aggregated URI in each page's file_task if that page was part of a successful aggregate.
            successful_files_query = """
                SELECT DISTINCT original_filename, gcs_output_markdown_uri 
                FROM file_tasks 
                WHERE batch_job_id = ? AND status = 'completed' AND gcs_output_markdown_uri IS NOT NULL
            """
            successful_file_uris_tuples = con.execute(
                successful_files_query, (task_id,)
            ).fetchall()

            if successful_file_uris_tuples:
                storage_client = storage.Client(
                    project=get_gcs_project_id(), credentials=get_gcs_credentials()
                )
                for original_fn, gcs_uri in successful_file_uris_tuples:
                    try:
                        bucket_name, blob_name = gcs_uri.replace("gs://", "").split(
                            "/", 1
                        )
                        bucket = storage_client.bucket(bucket_name)
                        blob = bucket.blob(blob_name)
                        markdown_content = await asyncio.to_thread(
                            blob.download_as_text
                        )
                        outputs_list.append(
                            BatchOutputItem(
                                original_filename=original_fn,
                                markdown_content=markdown_content,
                                gcs_output_uri=gcs_uri,
                            )
                        )
                    except Exception as e:
                        print(f"Error downloading GCS content for {gcs_uri}: {e}")
                        # If we can't download a supposedly successful file, list it as an error for this retrieval attempt
                        errors_list.append(
                            BatchJobFailedFileOutput(
                                original_filename=original_fn,
                                file_type="unknown_at_retrieval",  # We don't easily have file_type here from this query
                                error_message=f"Failed to download content from GCS: {str(e)}",
                                status="retrieval_error",
                            )
                        )

            if job_status == "completed_with_errors":
                message = "Job completed with some errors. See errors list."
                failed_tasks_query = """
                    SELECT original_filename, file_type, page_number, error_message, status 
                    FROM file_tasks 
                    WHERE batch_job_id = ? AND status = 'failed'
                """
                failed_tasks_tuples = con.execute(
                    failed_tasks_query, (task_id,)
                ).fetchall()
                for (
                    ft_orig_fn,
                    ft_type,
                    ft_page_num,
                    ft_err_msg,
                    ft_status,
                ) in failed_tasks_tuples:
                    errors_list.append(
                        BatchJobFailedFileOutput(
                            original_filename=ft_orig_fn,
                            file_type=ft_type,
                            page_number=ft_page_num,
                            error_message=ft_err_msg,
                            status=ft_status,
                        )
                    )
            elif (
                not outputs_list and not errors_list
            ):  # Status was 'completed' but no files/uris found or downloaded
                message = "Job reported as completed, but no output files were found or could be retrieved."

        elif job_status == "completed_no_files":
            message = "Job completed, but no files were processed (e.g., no supported files in input)."
        else:  # Pending, processing, or failed states where individual outputs are not expected
            message = f"Job is not yet fully completed or has failed. Current status: {job_status}"

        return BatchJobOutputResponse(
            job_id=job_id,
            status=job_status,
            outputs=outputs_list,
            errors=errors_list,
            message=message,
        )
    finally:
        con.close()


async def _process_single_non_pdf_file_and_upload(
    temp_file_path: str,
    file_ext: str,
    original_filename: str,
    output_gcs_path_str: str,
    main_batch_job_id: str,
    file_task_id: str,
):
    current_time = datetime.utcnow()
    con = get_db_connection()
    try:
        con.execute(
            "UPDATE file_tasks SET status = ?, updated_at = ? WHERE file_task_id = ?",
            ("processing", current_time, file_task_id),
        )
        con.execute(
            "UPDATE batch_jobs SET last_updated_at = ? WHERE job_id = ?",
            (current_time, main_batch_job_id),
        )
        con.commit()

        # Read content from temp file
        with open(temp_file_path, "rb") as f:
            content_bytes = f.read()

        # Re-use the existing _process_file_content logic
        markdown_texts = await _process_file_content(
            file_ext, content_bytes, get_pdf_backend()
        )

        if (
            not markdown_texts
            or markdown_texts[0].startswith("Error processing")
            or markdown_texts[0].startswith("Unsupported file type")
        ):
            error_message = (
                markdown_texts[0] if markdown_texts else "Unknown processing error"
            )
            con.execute(
                "UPDATE file_tasks SET status = ?, error_message = ?, updated_at = ? WHERE file_task_id = ?",
                ("failed", error_message, datetime.utcnow(), file_task_id),
            )
            con.execute(
                "UPDATE batch_jobs SET overall_failed_count = overall_failed_count + 1, last_updated_at = ? WHERE job_id = ?",
                (datetime.utcnow(), main_batch_job_id),
            )
            con.commit()
            print(
                f"Failed to process non-PDF file {original_filename}: {error_message}"
            )
            return

        # Join all markdown texts with a separator
        full_markdown_output = "\n\n---\n\n".join(markdown_texts)

        # Upload to GCS
        storage_client = storage.Client(
            project=get_gcs_project_id(), credentials=get_gcs_credentials()
        )

        # Parse the output GCS path
        if not output_gcs_path_str.startswith("gs://"):
            raise ValueError("Output GCS path must start with gs://")
        bucket_name = output_gcs_path_str.replace("gs://", "").split("/")[0]
        output_prefix = "/".join(
            output_gcs_path_str.replace("gs://", "").split("/")[1:]
        )
        if output_prefix and not output_prefix.endswith("/"):
            output_prefix += "/"

        # Create the output blob name
        base_name = os.path.splitext(original_filename)[0]
        output_blob_name = f"{output_prefix}{base_name}.md"
        bucket = storage_client.bucket(bucket_name)
        output_blob_obj = bucket.blob(output_blob_name)

        # Upload the markdown content
        await asyncio.to_thread(
            output_blob_obj.upload_from_string(
                full_markdown_output, content_type="text/markdown"
            )
        )
        gcs_output_url = f"gs://{bucket_name}/{output_blob_name}"

        # Update the task status in the database
        con.execute(
            "UPDATE file_tasks SET status = ?, gcs_output_markdown_uri = ?, updated_at = ? WHERE file_task_id = ?",
            ("completed", gcs_output_url, datetime.utcnow(), file_task_id),
        )
        con.execute(
            "UPDATE batch_jobs SET overall_processed_count = overall_processed_count + 1, last_updated_at = ? WHERE job_id = ?",
            (datetime.utcnow(), main_batch_job_id),
        )
        con.commit()

    except Exception as e:
        error_str = f"Error in _process_single_non_pdf_file_and_upload for {original_filename}: {str(e)}"
        print(error_str)
        try:
            # Attempt to mark as failed in DB
            con.execute(
                "UPDATE file_tasks SET status = ?, error_message = ?, updated_at = ? WHERE file_task_id = ?",
                ("failed", error_str, datetime.utcnow(), file_task_id),
            )
            con.execute(
                "UPDATE batch_jobs SET overall_failed_count = overall_failed_count + 1, last_updated_at = ? WHERE job_id = ?",
                (datetime.utcnow(), main_batch_job_id),
            )
            con.commit()
        except Exception as db_err:
            print(
                f"Additionally, failed to update DB for task {file_task_id} failure: {db_err}"
            )
    finally:
        # Clean up the temporary file
        try:
            os.unlink(temp_file_path)
        except Exception as cleanup_err:
            print(f"Error cleaning up temp file {temp_file_path}: {cleanup_err}")

        _check_and_finalize_batch_job_status(main_batch_job_id, con)
        con.commit()
        con.close()


def _run_gemini_pdf_batch_conversion(
    pdf_inputs_list: List[tuple[str, str]],  # List of (filename, temp_file_path)
    output_gcs_path_str: str,
    main_batch_job_id: str,
    gemini_batch_sub_job_id: str,
):
    """
    Process multiple PDFs using Gemini Batch Prediction.
    Now accepts a list of tuples containing (filename, temp_file_path).
    """
    current_time = datetime.utcnow()
    con = get_db_connection()
    storage_client = None
    gemini_client = None

    try:
        # Update status to preparing
        con.execute(
            "UPDATE gemini_pdf_batch_sub_jobs SET status = ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
            ("preparing", current_time, gemini_batch_sub_job_id),
        )
        con.commit()

        # Initialize clients
        storage_client = storage.Client(
            project=get_gcs_project_id(), credentials=get_gcs_credentials()
        )
        temp_bucket = storage_client.bucket(GCS_BATCH_BUCKET)
        gemini_client = get_gemini_client()

        # Process each PDF
        page_tasks_by_pdf: Dict[str, List[str]] = {}  # filename -> [task_ids]
        for original_pdf_filename, temp_file_path in pdf_inputs_list:
            try:
                # Convert PDF to images
                with open(temp_file_path, "rb") as pdf_file:
                    page_images = convert_from_bytes(pdf_file.read(), fmt="png")
            except Exception as e:
                print(f"Failed to convert PDF {original_pdf_filename} to images: {e}")
                error_message = f"Failed to convert PDF to images: {str(e)}"
                _record_pdf_failure(
                    con,
                    main_batch_job_id,
                    original_pdf_filename,
                    error_message,
                    current_time,
                )
                continue

            # Create task records and upload page images
            page_task_ids = []
            for page_num, page_image in enumerate(page_images, start=1):
                try:
                    # Create a task record for this page
                    page_task_id = str(uuid.uuid4())
                    page_task_ids.append(page_task_id)

                    con.execute(
                        """INSERT INTO file_tasks 
                        (file_task_id, batch_job_id, gemini_batch_sub_job_id, original_filename, file_type, page_number, status, created_at, updated_at) 
                        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)""",
                        (
                            page_task_id,
                            main_batch_job_id,
                            gemini_batch_sub_job_id,
                            original_pdf_filename,
                            ".pdf",
                            page_num,
                            "pending",
                            current_time,
                            current_time,
                        ),
                    )

                    # Save page image to temporary GCS location
                    temp_image_name = f"temp/{main_batch_job_id}/{page_task_id}.png"
                    temp_blob = temp_bucket.blob(temp_image_name)

                    # Save the page image to a bytes buffer and upload
                    img_byte_arr = BytesIO()
                    page_image.save(img_byte_arr, format="PNG")
                    img_byte_arr.seek(0)
                    temp_blob.upload_from_file(img_byte_arr, content_type="image/png")
                    image_gcs_uri = f"gs://{GCS_BATCH_BUCKET}/{temp_image_name}"

                    # Update task with image URI
                    con.execute(
                        "UPDATE file_tasks SET gcs_input_image_uri = ?, status = ?, updated_at = ? WHERE file_task_id = ?",
                        (image_gcs_uri, "image_uploaded", current_time, page_task_id),
                    )

                except Exception as e:
                    print(
                        f"Error processing page {page_num} of {original_pdf_filename}: {e}"
                    )
                    error_message = f"Failed to process page {page_num}: {str(e)}"
                    con.execute(
                        "UPDATE file_tasks SET status = ?, error_message = ?, updated_at = ? WHERE file_task_id = ?",
                        ("failed", error_message, current_time, page_task_id),
                    )
                    continue

            if page_task_ids:  # If we successfully processed any pages
                page_tasks_by_pdf[original_pdf_filename] = page_task_ids
            con.commit()

            # Clean up the temporary PDF file
            try:
                os.unlink(temp_file_path)
            except Exception as cleanup_err:
                print(f"Error cleaning up temp file {temp_file_path}: {cleanup_err}")

        if not page_tasks_by_pdf:
            error_message = (
                "No PDF pages were successfully prepared for batch processing"
            )
            con.execute(
                "UPDATE gemini_pdf_batch_sub_jobs SET status = ?, error_message = ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
                ("failed", error_message, current_time, gemini_batch_sub_job_id),
            )
            con.commit()
            return

        # Prepare Gemini batch job
        temp_gcs_input_prefix = (
            f"gemini_batch_jobs/{main_batch_job_id}/{gemini_batch_sub_job_id}/inputs"
        )
        temp_gcs_output_prefix = (
            f"gemini_batch_jobs/{main_batch_job_id}/{gemini_batch_sub_job_id}/outputs"
        )

        # Create payload for Gemini batch job
        payload_items = []
        for task_info in con.execute(
            "SELECT file_task_id, gcs_input_image_uri FROM file_tasks WHERE gemini_batch_sub_job_id = ? AND status = 'image_uploaded'",
            (gemini_batch_sub_job_id,),
        ).fetchall():
            task_id, image_uri = task_info
            payload_items.append(
                {
                    "id": task_id,
                    "request": {
                        "contents": [
                            {
                                "role": "user",
                                "parts": [
                                    {
                                        "file_data": {
                                            "file_uri": image_uri,
                                            "mime_type": "image/png",
                                        }
                                    },
                                    {"text": OCR_PROMPT},
                                ],
                            }
                        ]
                    },
                }
            )

        # Upload payload.jsonl
        payload_jsonl = "\n".join(json.dumps(item) for item in payload_items)
        payload_blob_name = f"{temp_gcs_input_prefix}/payload.jsonl"
        payload_blob = temp_bucket.blob(payload_blob_name)
        payload_blob.upload_from_string(payload_jsonl, content_type="application/jsonl")
        payload_gcs_uri = f"gs://{GCS_BATCH_BUCKET}/{payload_blob_name}"

        # Update status
        con.execute(
            "UPDATE gemini_pdf_batch_sub_jobs SET status = ?, payload_gcs_uri = ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
            ("submitting", current_time, payload_gcs_uri, gemini_batch_sub_job_id),
        )
        con.commit()

        # Submit Gemini batch job
        gemini_output_uri = f"gs://{GCS_BATCH_BUCKET}/{temp_gcs_output_prefix}"
        batch_job_config = CreateBatchJobConfig(dest=gemini_output_uri)

        gemini_job = gemini_client.batches.create(
            model=GEMINI_MODEL_FOR_VISION,
            src=payload_gcs_uri,
            config=batch_job_config,
        )

        # Update with Gemini job details
        con.execute(
            """UPDATE gemini_pdf_batch_sub_jobs 
            SET status = ?, gemini_api_job_name = ?, gemini_output_gcs_uri_prefix = ?, updated_at = ? 
            WHERE gemini_batch_sub_job_id = ?""",
            (
                str(gemini_job.state),
                gemini_job.name,
                gemini_output_uri,
                current_time,
                gemini_batch_sub_job_id,
            ),
        )
        con.commit()

        # Process results if job completed successfully
        if gemini_job.state == JobState.JOB_STATE_SUCCEEDED:
            # Find and process predictions file
            predictions_blob = None
            for blob in storage_client.list_blobs(
                GCS_BATCH_BUCKET, prefix=f"{temp_gcs_output_prefix}/"
            ):
                if "predictions.jsonl" in blob.name:
                    predictions_blob = blob
                    break

            if predictions_blob:
                predictions_content = predictions_blob.download_as_text()

                # Process each prediction and update tasks
                for line in predictions_content.splitlines():
                    if not line.strip():
                        continue

                    prediction = json.loads(line)
                    task_id = prediction["id"]
                    markdown_text = prediction["response"]["candidates"][0]["content"][
                        "parts"
                    ][0]["text"]

                    # Update task with result
                    con.execute(
                        """UPDATE file_tasks 
                        SET status = ?, markdown_content = ?, updated_at = ? 
                        WHERE file_task_id = ?""",
                        ("completed", markdown_text, current_time, task_id),
                    )

                # Aggregate results by PDF and upload final markdown files
                for pdf_filename, task_ids in page_tasks_by_pdf.items():
                    # Get all successful pages for this PDF
                    pages_data = con.execute(
                        """SELECT page_number, markdown_content 
                        FROM file_tasks 
                        WHERE file_task_id IN (SELECT UNNEST(?)) 
                        AND status = 'completed'
                        ORDER BY page_number""",
                        (task_ids,),
                    ).fetchall()

                    if pages_data:
                        # Combine all pages
                        full_markdown = "\n\n---\n\n".join(
                            page[1] for page in pages_data
                        )

                        # Upload final markdown
                        output_bucket_name = output_gcs_path_str.replace(
                            "gs://", ""
                        ).split("/")[0]
                        output_prefix = "/".join(
                            output_gcs_path_str.replace("gs://", "").split("/")[1:]
                        )
                        if output_prefix and not output_prefix.endswith("/"):
                            output_prefix += "/"

                        final_blob_name = (
                            f"{output_prefix}{os.path.splitext(pdf_filename)[0]}.md"
                        )
                        final_bucket = storage_client.bucket(output_bucket_name)
                        final_blob = final_bucket.blob(final_blob_name)
                        final_blob.upload_from_string(
                            full_markdown, content_type="text/markdown"
                        )

                        # Update all tasks for this PDF with the final URI
                        final_uri = f"gs://{output_bucket_name}/{final_blob_name}"
                        con.execute(
                            """UPDATE file_tasks 
                            SET gcs_output_markdown_uri = ? 
                            WHERE file_task_id IN (SELECT UNNEST(?))""",
                            (final_uri, task_ids),
                        )

                # Update sub-job status
                con.execute(
                    """UPDATE gemini_pdf_batch_sub_jobs 
                    SET status = 'completed', updated_at = ? 
                    WHERE gemini_batch_sub_job_id = ?""",
                    (current_time, gemini_batch_sub_job_id),
                )
                con.commit()

            else:
                raise Exception("Predictions file not found in Gemini output")

        else:
            error_message = f"Gemini job failed with state: {gemini_job.state}"
            raise Exception(error_message)

    except Exception as e:
        error_message = f"Error in PDF batch conversion: {str(e)}"
        print(error_message)
        try:
            con.execute(
                "UPDATE gemini_pdf_batch_sub_jobs SET status = ?, error_message = ?, updated_at = ? WHERE gemini_batch_sub_job_id = ?",
                ("failed", error_message, current_time, gemini_batch_sub_job_id),
            )
            con.commit()
        except Exception as db_err:
            print(f"Additionally, failed to update DB for failure: {db_err}")
    finally:
        if con:
            _check_and_finalize_batch_job_status(main_batch_job_id, con)
            con.commit()
            con.close()


def _check_and_finalize_batch_job_status(
    main_batch_job_id: str, con: duckdb.DuckDBPyConnection
):
    """Checks if all tasks for a batch job are completed and updates the main job status.
    This function assumes the connection `con` is open and does not close it.
    """
    try:
        job_info = con.execute(
            "SELECT total_input_files, overall_processed_count, overall_failed_count, status FROM batch_jobs WHERE job_id = ?",
            (main_batch_job_id,),
        ).fetchone()

        if not job_info:
            print(
                f"_check_and_finalize_batch_job_status: Batch job {main_batch_job_id} not found."
            )
            return

        total_files, processed_count, failed_count, current_status = job_info

        # If already in a final state, no need to update further by this check.
        if current_status in [
            "completed",
            "completed_with_errors",
            "failed_catastrophic",
            "completed_no_files",
        ]:
            return

        if (processed_count + failed_count) >= total_files:
            new_status = ""
            if failed_count > 0:
                new_status = "completed_with_errors"
            else:
                new_status = "completed"

            if new_status:
                print(
                    f"Finalizing batch job {main_batch_job_id} to status: {new_status}"
                )
                con.execute(
                    "UPDATE batch_jobs SET status = ?, last_updated_at = ? WHERE job_id = ?",
                    (new_status, datetime.utcnow(), main_batch_job_id),
                )
    except Exception as e:
        print(
            f"Error in _check_and_finalize_batch_job_status for {main_batch_job_id}: {e}"
        )
        # Optionally, update main batch job to a specific error state if this check itself fails critically


def _record_pdf_failure(
    con: duckdb.DuckDBPyConnection,
    batch_job_id: str,
    pdf_filename: str,
    error_message: str,
    current_time: datetime,
):
    """Helper function to record PDF failures in the database."""
    file_task_id = str(uuid.uuid4())
    con.execute(
        """INSERT INTO file_tasks 
        (file_task_id, batch_job_id, original_filename, file_type, status, error_message, created_at, updated_at) 
        VALUES (?, ?, ?, ?, ?, ?, ?, ?)""",
        (
            file_task_id,
            batch_job_id,
            pdf_filename,
            ".pdf",
            "failed",
            error_message,
            current_time,
            current_time,
        ),
    )
    con.execute(
        "UPDATE batch_jobs SET overall_failed_count = overall_failed_count + 1, last_updated_at = ? WHERE job_id = ?",
        (current_time, batch_job_id),
    )
    con.commit()


# --- Main function for Uvicorn ---
def main():
    import uvicorn

    # Read host and port from environment variables or use defaults
    host = os.getenv("LLM_FOOD_HOST", "0.0.0.0")
    port = int(os.getenv("LLM_FOOD_PORT", "8000"))
    reload = (
        os.getenv("LLM_FOOD_RELOAD", "false").lower() == "true"
    )  # Added reload option

    print(
        f"Starting server on {host}:{port} with reload={'enabled' if reload else 'disabled'}"
    )
    uvicorn.run(
        "llm_food.app:app", host=host, port=port, reload=reload
    )  # Corrected to pass app string for reload


if __name__ == "__main__":
    # This allows running the FastAPI app directly using `python -m llm_food.app`
    main()
